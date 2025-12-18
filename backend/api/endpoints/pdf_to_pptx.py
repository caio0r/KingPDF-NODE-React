from fastapi import APIRouter, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
import shutil
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from core.utils import cleanup_file
from PIL import Image
import io

router = APIRouter()

UPLOAD_DIR = "uploads"
OUTPUT_DIR = "outputs"

@router.post("/pdf-to-pptx")
def pdf_to_pptx(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Invalid file type. Please upload a PDF.")

    input_path = os.path.join(UPLOAD_DIR, f"pptx_in_{file.filename}")
    output_filename = f"{os.path.splitext(file.filename)[0]}.pptx"
    output_path = os.path.join(OUTPUT_DIR, output_filename)

    try:
        # Save uploaded file
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Open PDF with PyMuPDF
        pdf_document = fitz.open(input_path)
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Convert each PDF page to an image and add to PowerPoint
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Render page to image (matrix for higher resolution)
            mat = fitz.Matrix(2, 2)  # 2x zoom for better quality
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PIL Image
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Save temporary image
            temp_img_path = os.path.join(OUTPUT_DIR, f"temp_slide_{page_num}.png")
            img.save(temp_img_path)
            
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide (fill the entire slide)
            slide.shapes.add_picture(
                temp_img_path,
                0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # Clean up temporary image
            cleanup_file(temp_img_path)

        pdf_document.close()

        # Save PowerPoint
        prs.save(output_path)

        background_tasks.add_task(cleanup_file, output_path)

        return FileResponse(
            output_path, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            filename=output_filename
        )

    except Exception as e:
        print(f"Conversion error: {e}")
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")
    
    finally:
        cleanup_file(input_path)
