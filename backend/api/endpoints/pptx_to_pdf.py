from fastapi import APIRouter, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
import shutil
import os
from pptx import Presentation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Image as RLImage, PageBreak
from reportlab.lib.units import inch
from core.utils import cleanup_file
from PIL import Image
import io

router = APIRouter()

UPLOAD_DIR = "uploads"
OUTPUT_DIR = "outputs"

@router.post("/pptx-to-pdf")
async def pptx_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    if not (file.filename.endswith(".pptx") or file.filename.endswith(".ppt")):
        raise HTTPException(status_code=400, detail="Invalid file type. Please upload a PowerPoint file (.pptx or .ppt).")

    # Note: .ppt files (old format) are not fully supported, primarily .pptx
    if file.filename.endswith(".ppt"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported. Please convert your .ppt file to .pptx first.")

    input_path = os.path.join(UPLOAD_DIR, f"pptx_in_{file.filename}")
    output_filename = f"{os.path.splitext(file.filename)[0]}.pdf"
    output_path = os.path.join(OUTPUT_DIR, output_filename)

    try:
        # Save uploaded file
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Open PowerPoint presentation
        prs = Presentation(input_path)
        
        # Create PDF with landscape orientation (typical for presentations)
        pdf = SimpleDocTemplate(
            output_path, 
            pagesize=(11*inch, 8.5*inch),  # Landscape letter size
            rightMargin=0.5*inch, 
            leftMargin=0.5*inch,
            topMargin=0.5*inch, 
            bottomMargin=0.5*inch
        )
        
        # Container for the 'Flowable' objects
        elements = []
        
        slide_count = 0
        
        # Process each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_count += 1
            
            # Create a simple representation of the slide
            # For now, we'll extract text and images from each slide
            
            # Extract text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # If there's text, add it to PDF
            if slide_text:
                from reportlab.platypus import Paragraph
                from reportlab.lib.styles import getSampleStyleSheet
                
                styles = getSampleStyleSheet()
                title_style = styles['Heading1']
                normal_style = styles['Normal']
                
                # First line as title, rest as content
                if len(slide_text) > 0:
                    title_para = Paragraph(slide_text[0], title_style)
                    elements.append(title_para)
                    
                    from reportlab.platypus import Spacer
                    elements.append(Spacer(1, 12))
                    
                    for text in slide_text[1:]:
                        para = Paragraph(text, normal_style)
                        elements.append(para)
                        elements.append(Spacer(1, 6))
            
            # Extract images from slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Save temp image
                        temp_img_path = os.path.join(OUTPUT_DIR, f"temp_slide_{slide_idx}.png")
                        with open(temp_img_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                        
                        # Open with PIL to get dimensions
                        pil_img = Image.open(temp_img_path)
                        img_width, img_height = pil_img.size
                        
                        # Calculate scaled dimensions
                        max_width = 9 * inch
                        max_height = 6.5 * inch
                        
                        aspect = img_height / float(img_width)
                        
                        if img_width > max_width:
                            img_width = max_width
                            img_height = img_width * aspect
                        
                        if img_height > max_height:
                            img_height = max_height
                            img_width = img_height / aspect
                        
                        # Add image to PDF
                        img = RLImage(temp_img_path, width=img_width, height=img_height)
                        elements.append(img)
                        
                        # Clean up
                        cleanup_file(temp_img_path)
                    except Exception as e:
                        print(f"Error processing image in slide {slide_idx}: {e}")
            
            # Add page break between slides (except for the last one)
            if slide_idx < len(prs.slides) - 1:
                elements.append(PageBreak())
        
        # Build PDF
        if elements:
            pdf.build(elements)
            print(f"Successfully converted {slide_count} slides to PDF")
        else:
            # Create empty message
            from reportlab.platypus import Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            styles = getSampleStyleSheet()
            elements.append(Paragraph("A apresentação não contém conteúdo visível.", styles['Normal']))
            pdf.build(elements)

        if not os.path.exists(output_path):
            raise HTTPException(status_code=500, detail="Conversion failed: Output file not created.")

        background_tasks.add_task(cleanup_file, output_path)

        return FileResponse(
            output_path, 
            media_type="application/pdf", 
            filename=output_filename
        )

    except Exception as e:
        print(f"Conversion error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")
    
    finally:
        cleanup_file(input_path)
